import os
from pathlib import Path
from typing import Any, Dict, List, Optional
from .base import BaseTool, ToolContext, ToolResult
from src.security.guard import ActionRiskLevel
from src.security.snapshot import get_snapshot_manager

def resolve_target_path(raw_path: str) -> Path:
    p = os.path.expanduser(str(raw_path).strip())
    p_lower = p.lower()
    if p_lower.startswith("desktop/") or p_lower.startswith("desktop\\"):
        return (Path.home() / "Desktop" / p[8:]).resolve()
    elif p_lower.startswith("documents/") or p_lower.startswith("documents\\"):
        return (Path.home() / "Documents" / p[10:]).resolve()
    elif p_lower.startswith("downloads/") or p_lower.startswith("downloads\\"):
        return (Path.home() / "Downloads" / p[10:]).resolve()
    return Path(p).resolve()

class GenerateDocumentTool(BaseTool):
    name = "generate_document"
    description = (
        "Create rich, professionally formatted documents on the local filesystem. "
        "Supports 'docx' (Word), 'xlsx' (Excel spreadsheet), 'pdf' (PDF document), "
        "'pptx' (PowerPoint presentation), and 'md' (Markdown)."
    )
    parameters = {
        "type": "object",
        "properties": {
            "file_path": {"type": "string", "description": "Destination file path (e.g. 'C:/Users/.../report.docx' or 'Desktop/report.docx')."},
            "doc_type": {
                "type": "string",
                "enum": ["docx", "xlsx", "pdf", "pptx", "md"],
                "description": "Document format to generate."
            },
            "title": {"type": "string", "description": "Document title."},
            "sections": {
                "type": "array",
                "description": "Structured sections of the document. Each section has 'heading', 'content', and optional 'table_data' (list of rows) or 'bullet_points' (list of strings).",
                "items": {
                    "type": "object",
                    "properties": {
                        "heading": {"type": "string"},
                        "content": {"type": "string"},
                        "bullet_points": {"type": "array", "items": {"type": "string"}},
                        "table_data": {"type": "array", "items": {"type": "array", "items": {"type": "string"}}}
                    },
                    "required": ["heading"]
                }
            }
        },
        "required": ["file_path", "doc_type", "title", "sections"]
    }
    risk_level = ActionRiskLevel.MODERATE

    async def execute(self, args: Dict[str, Any], context: ToolContext) -> ToolResult:
        file_path = resolve_target_path(str(args.get("file_path")))
        doc_type = str(args.get("doc_type", "docx")).lower()
        title = str(args.get("title", "Document"))
        sections = args.get("sections", [])

        # Snapshot destination if already exists
        snapshot_mgr = get_snapshot_manager()
        snap_id = snapshot_mgr.snapshot_before_action(
            target_path=file_path,
            action_type="modify" if file_path.exists() else "create",
            task_id=context.task_id,
            description=f"Generating {doc_type.upper()} document: {file_path.name}"
        )

        try:
            file_path.parent.mkdir(parents=True, exist_ok=True)

            if doc_type == "docx":
                self._generate_docx(file_path, title, sections)
            elif doc_type == "xlsx":
                self._generate_xlsx(file_path, title, sections)
            elif doc_type == "pdf":
                self._generate_pdf(file_path, title, sections)
            elif doc_type == "pptx":
                self._generate_pptx(file_path, title, sections)
            elif doc_type == "md":
                self._generate_md(file_path, title, sections)
            else:
                return ToolResult(success=False, data=None, error=f"Unsupported document format: {doc_type}")

            return ToolResult(
                success=True,
                data={
                    "file_path": str(file_path),
                    "doc_type": doc_type,
                    "title": title,
                    "sections_count": len(sections),
                    "snapshot_id": snap_id,
                    "status": "Document created successfully"
                },
                snapshot_id=snap_id
            )
        except Exception as e:
            return ToolResult(success=False, data=None, error=f"Failed to generate document: {str(e)}")

    def _generate_docx(self, path: Path, title: str, sections: List[Dict[str, Any]]):
        import docx
        from docx.shared import Inches, Pt, RGBColor

        doc = docx.Document()
        doc.add_heading(title, 0)

        for sec in sections:
            heading = sec.get("heading")
            if heading:
                doc.add_heading(heading, level=1)
            content = sec.get("content")
            if content:
                doc.add_paragraph(content)

            bullets = sec.get("bullet_points", [])
            for bp in bullets:
                doc.add_paragraph(bp, style="List Bullet")

            table_data = sec.get("table_data", [])
            if table_data and len(table_data) > 0:
                rows = len(table_data)
                cols = max(len(r) for r in table_data)
                table = doc.add_table(rows=rows, cols=cols)
                table.style = "Light Shading Accent 1" if "Light Shading Accent 1" in [s.name for s in doc.styles] else "Table Grid"
                for r_idx, row in enumerate(table_data):
                    for c_idx, val in enumerate(row):
                        table.cell(r_idx, c_idx).text = str(val)

        doc.save(str(path))

    def _generate_xlsx(self, path: Path, title: str, sections: List[Dict[str, Any]]):
        import openpyxl
        from openpyxl.styles import Font, PatternFill, Alignment
        import re

        wb = openpyxl.Workbook()
        ws = wb.active
        clean_title = re.sub(r'[\\/*?:\[\]]', '_', title)[:31] if title else "Sheet1"
        ws.title = clean_title or "Sheet1"

        current_row = 1
        ws.cell(row=current_row, column=1, value=title).font = Font(size=16, bold=True)
        current_row += 2

        for sec in sections:
            heading = sec.get("heading")
            if heading:
                cell = ws.cell(row=current_row, column=1, value=heading)
                cell.font = Font(size=12, bold=True, color="FFFFFF")
                cell.fill = PatternFill(start_color="1F497D", end_color="1F497D", fill_type="solid")
                current_row += 1

            content = sec.get("content")
            if content:
                ws.cell(row=current_row, column=1, value=content)
                current_row += 1

            table_data = sec.get("table_data", [])
            if table_data:
                for r_idx, row in enumerate(table_data):
                    for c_idx, val in enumerate(row):
                        cell = ws.cell(row=current_row, column=c_idx + 1, value=val)
                        if r_idx == 0:
                            cell.font = Font(bold=True)
                    current_row += 1

            current_row += 1

        wb.save(str(path))

    def _generate_pdf(self, path: Path, title: str, sections: List[Dict[str, Any]]):
        from reportlab.lib.pagesizes import letter
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib import colors
        import xml.sax.saxutils

        def safe_xml(s: Any) -> str:
            return xml.sax.saxutils.escape(str(s))

        doc = SimpleDocTemplate(str(path), pagesize=letter)
        styles = getSampleStyleSheet()
        story = []

        story.append(Paragraph(safe_xml(title), styles['Title']))
        story.append(Spacer(1, 14))

        for sec in sections:
            heading = sec.get("heading")
            if heading:
                story.append(Paragraph(safe_xml(heading), styles['Heading2']))
                story.append(Spacer(1, 8))

            content = sec.get("content")
            if content:
                story.append(Paragraph(safe_xml(content), styles['Normal']))
                story.append(Spacer(1, 8))

            bullets = sec.get("bullet_points", [])
            for bp in bullets:
                story.append(Paragraph(f"• {safe_xml(bp)}", styles['Normal']))
                story.append(Spacer(1, 4))

            table_data = sec.get("table_data", [])
            if table_data:
                # Stringify table cells
                cleaned_table = [[str(c) for c in row] for row in table_data]
                t = Table(cleaned_table)
                t.setStyle(TableStyle([
                    ('BACKGROUND', (0, 0), (-1, 0), colors.navy),
                    ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                    ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
                    ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                    ('BOTTOMPADDING', (0, 0), (-1, 0), 6),
                    ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
                    ('GRID', (0, 0), (-1, -1), 1, colors.black),
                ]))
                story.append(t)
                story.append(Spacer(1, 12))

        doc.build(story)

    def _generate_pptx(self, path: Path, title: str, sections: List[Dict[str, Any]]):
        from pptx import Presentation
        from pptx.util import Inches, Pt

        prs = Presentation()
        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        slide.shapes.title.text = title

        # Section slides
        bullet_slide_layout = prs.slide_layouts[1]
        for sec in sections:
            slide = prs.slides.add_slide(bullet_slide_layout)
            slide.shapes.title.text = sec.get("heading", "Slide")
            tf = slide.placeholders[1].text_frame

            content = sec.get("content")
            if content:
                tf.text = content

            bullets = sec.get("bullet_points", [])
            for bp in bullets:
                p = tf.add_paragraph()
                p.text = bp
                p.level = 0

        prs.save(str(path))

    def _generate_md(self, path: Path, title: str, sections: List[Dict[str, Any]]):
        lines = [f"# {title}\n\n"]
        for sec in sections:
            heading = sec.get("heading")
            if heading:
                lines.append(f"## {heading}\n\n")
            content = sec.get("content")
            if content:
                lines.append(f"{content}\n\n")
            bullets = sec.get("bullet_points", [])
            for bp in bullets:
                lines.append(f"- {bp}\n")
            if bullets:
                lines.append("\n")
            table_data = sec.get("table_data", [])
            if table_data and len(table_data) > 0:
                header = table_data[0]
                lines.append("| " + " | ".join(str(c) for c in header) + " |\n")
                lines.append("| " + " | ".join(["---"] * len(header)) + " |\n")
                for row in table_data[1:]:
                    lines.append("| " + " | ".join(str(c) for c in row) + " |\n")
                lines.append("\n")

        with open(path, "w", encoding="utf-8") as f:
            f.writelines(lines)
